"""PPTX generation from template files via python-pptx."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

from app.sandbox.path_guard import normalize_user_path
from app.tools.builtin.docgen.schemas import PptxSlide, TemplateId

_TEMPLATE_IDS: tuple[TemplateId, ...] = ("minimal", "business", "creative")


def templates_dir() -> Path:
    """Return ``skills/_templates`` relative to the repo root."""
    # backend/app/tools/builtin/docgen/pptx_gen.py → repo root is 5 parents up
    return Path(__file__).resolve().parents[5] / "skills" / "_templates"


def template_path(template_id: TemplateId | str) -> Path:
    if template_id not in _TEMPLATE_IDS:
        raise ValueError(f"Unknown template_id: {template_id!r}. Available: {_TEMPLATE_IDS}")
    path = templates_dir() / f"tpl-{template_id}.pptx"
    if not path.is_file():
        raise FileNotFoundError(f"Template not found: {path}")
    return path


def _set_slide_content(slide, title: str, bullets: list[str]) -> None:
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    # Prefer body placeholder; fall back to first non-title text frame.
    body = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape == slide.shapes.title:
            continue
        body = shape
        break
    if body is None:
        return
    tf = body.text_frame
    tf.clear()
    if not bullets:
        return
    p0 = tf.paragraphs[0]
    p0.text = bullets[0]
    p0.level = 0
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)


async def gen(
    template_id: TemplateId | str,
    slides: list[PptxSlide] | list[dict],
    out_path: str,
) -> str:
    """Fill *template_id* with *slides* and save to *out_path*."""
    parsed = [
        s if isinstance(s, PptxSlide) else PptxSlide.model_validate(s) for s in slides
    ]
    target = normalize_user_path(out_path)
    target.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(template_path(template_id)))
    # Ensure we have enough slides (templates ship with 1 blank content slide).
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    while len(prs.slides) < len(parsed):
        prs.slides.add_slide(layout)
    # Remove extras if template had more than needed
    # python-pptx cannot easily delete slides; only fill what we need.
    for i, slide_data in enumerate(parsed):
        _set_slide_content(prs.slides[i], slide_data.title, slide_data.bullets)

    prs.save(str(target))
    return str(target)


def ensure_templates() -> None:
    """Create the three template pptx files if missing (dev bootstrap)."""
    dest = templates_dir()
    dest.mkdir(parents=True, exist_ok=True)
    for tid in _TEMPLATE_IDS:
        path = dest / f"tpl-{tid}.pptx"
        if path.is_file():
            continue
        prs = Presentation()
        # Title + content layout (index 1 if present, else title-only).
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = f"{{{{title}}}} ({tid})"
        # Add a body placeholder so _set_slide_content has somewhere to write
        # bullets. python-pptx does not expose a public add_placeholder API,
        # so we add a text box anchored to the body region.
        body = slide.shapes.add_textbox(
            left=prs.slide_width // 8,
            top=prs.slide_height // 4,
            width=prs.slide_width * 3 // 4,
            height=prs.slide_height * 2 // 3,
        )
        body.text_frame.text = "{{bullets}}"
        prs.save(str(path))
