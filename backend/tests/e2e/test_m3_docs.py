"""M3 document generation end-to-end smokes (fake LLM, real file writers)."""

from __future__ import annotations

from pathlib import Path

import httpx
import pytest
import respx
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from app.main import build_stack, create_app
from app.tools.builtin.docgen import docx_gen, pdf_gen, pptx_gen, xlsx_gen
from tests.conftest import FakeChatModel, make_ai


@pytest.mark.asyncio
async def test_m3_docx_smoke(tmp_path: Path):
    out = tmp_path / "brief.docx"
    path = await docx_gen.gen(
        {
            "title": "合同摘要",
            "paragraphs": [
                {"heading": "甲方", "body": "Acme"},
                {"heading": "标的", "body": "软件许可"},
            ],
        },
        str(out),
    )
    doc = Document(path)
    assert any(p.text == "合同摘要" for p in doc.paragraphs)


@pytest.mark.asyncio
async def test_m3_pptx_smoke(tmp_path: Path):
    pptx_gen.ensure_templates()
    out = tmp_path / "2026Q1.pptx"
    path = await pptx_gen.gen(
        "minimal",
        [
            {"title": "2026Q1", "bullets": ["目标"]},
            {"title": "进展", "bullets": ["完成 A", "完成 B"]},
            {"title": "风险", "bullets": ["延期风险"]},
            {"title": "计划", "bullets": ["Q2 重点"]},
            {"title": "总结", "bullets": ["继续推进"]},
        ],
        str(out),
    )
    prs = Presentation(path)
    assert len(prs.slides) >= 5
    assert prs.slides[0].shapes.title.text == "2026Q1"


@pytest.mark.asyncio
async def test_m3_xlsx_smoke(tmp_path: Path):
    out = tmp_path / "metrics.xlsx"
    path = await xlsx_gen.gen(
        {
            "headers": ["月", "收入", "成本"],
            "rows": [["1", 10, 4], ["2", 12, 5]],
            "sheet_name": "Q1",
        },
        str(out),
    )
    ws = load_workbook(path)["Q1"]
    assert ws.max_row == 3
    assert ws.max_column == 3


@pytest.mark.asyncio
@respx.mock
async def test_m3_pdf_smoke(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    monkeypatch.setenv("ELECTRON_PDF_PORT", "19876")
    out = tmp_path / "page.pdf"
    respx.post("http://127.0.0.1:19876/print-to-pdf").mock(
        return_value=httpx.Response(200, content=b"%PDF-1.4 m3")
    )
    path = await pdf_gen.gen("<h1>Q1</h1>", str(out))
    assert Path(path).read_bytes().startswith(b"%PDF")


@pytest.mark.asyncio
async def test_m3_pptx_via_chat_route(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    """Workforce → document_agent → pptx_gen with confirm auto-approved."""
    monkeypatch.setattr(
        "app.runtime.v2.office.officecli_available",
        lambda: False,
    )
    pptx_gen.ensure_templates()
    out = tmp_path / "via-chat.pptx"
    desk = str(out)

    planner = FakeChatModel(
        responses=[
            make_ai(
                content=(
                    '[{"id":"task_1","content":"生成 PPT《2026Q1》2 页",'
                    '"assignee":"document_agent","dependencies":[]}]'
                )
            ),
        ]
    )
    document = FakeChatModel(
        responses=[
            make_ai(
                content="",
                tool_calls=[
                    {
                        "name": "pptx_gen",
                        "args": {
                            "template_id": "business",
                            "slides": [
                                {"title": "封面", "bullets": ["2026Q1"]},
                                {"title": "内容", "bullets": ["要点"]},
                            ],
                            "out_path": desk,
                        },
                        "id": "call_pptx_1",
                    }
                ],
            ),
            make_ai(content="PPT 已生成。"),
        ]
    )
    idle = FakeChatModel(responses=[make_ai(content="ok")])

    stack = build_stack(
        supervisor_llm=planner,
        file_worker_llm=idle,
        doc_worker_llm=document,
        web_worker_llm=idle,
        msg_worker_llm=idle,
        whitelist=[str(tmp_path)],
    )
    task_manager = stack["task_manager"]
    bus = stack["bus"]
    confirm_hub = stack["confirm_hub"]

    def _auto_approve(event: dict) -> None:
        if event.get("type") == "tool.confirm_request":
            confirm_hub.resolve(event["call_id"], True)
        if event.get("type") == "to_sub_tasks":
            confirm_hub.resolve_plan(
                event.get("task_id") or "",
                event.get("subtasks") or [],
            )

    bus.subscribe(_auto_approve)
    app = create_app(task_manager=task_manager, bus=bus, confirm_hub=confirm_hub)

    async with httpx.AsyncClient(
        transport=httpx.ASGITransport(app=app), base_url="http://test"
    ) as client:
        async with client.stream(
            "POST", "/api/chat", json={"text": "生成 PPT《2026Q1》2 页"}
        ) as response:
            body = "".join([chunk async for chunk in response.aiter_text()])

    assert "graph.end" in body
    assert Path(desk).is_file()
    assert Presentation(desk).slides[0].shapes.title.text == "封面"
