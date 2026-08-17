"""Unit tests for pptx_gen."""

from pathlib import Path

import pytest
from pptx import Presentation

from app.tools.builtin.docgen import pptx_gen
from app.tools.builtin.docgen.coerce import coerce_pptx_slides
from app.tools.builtin.docgen.tools import PptxArgs


@pytest.fixture(scope="module", autouse=True)
def _templates():
    pptx_gen.ensure_templates()


@pytest.mark.asyncio
@pytest.mark.parametrize("template_id", ["minimal", "business", "creative"])
async def test_pptx_gen_fills_titles_and_bullets(tmp_path: Path, template_id: str):
    out = tmp_path / f"{template_id}.pptx"
    slides = [
        {"title": "封面", "bullets": ["副标题"]},
        {"title": "进展", "bullets": ["完成 A", "完成 B"]},
    ]
    path = await pptx_gen.gen(template_id, slides, str(out))
    prs = Presentation(path)
    assert len(prs.slides) >= 2
    assert prs.slides[0].shapes.title.text == "封面"
    assert prs.slides[1].shapes.title.text == "进展"
    body = None
    for shape in prs.slides[1].shapes:
        if shape.has_text_frame and shape != prs.slides[1].shapes.title:
            body = shape
            break
    assert body is not None
    texts = [p.text for p in body.text_frame.paragraphs if p.text.strip()]
    assert texts == ["完成 A", "完成 B"]


def test_coerce_slides_accepts_json_string_and_string_bullets():
    slides = coerce_pptx_slides(
        '[{"title":"封面","bullets":"副标题"},{"title":"行程","bullets":["D1","D2"]}]'
    )
    assert len(slides) == 2
    assert slides[0].title == "封面"
    assert slides[0].bullets == ["副标题"]
    assert slides[1].bullets == ["D1", "D2"]


def test_coerce_slides_accepts_single_object():
    slides = coerce_pptx_slides({"title": "Only", "bullet_points": ["a", "b"]})
    assert len(slides) == 1
    assert slides[0].bullets == ["a", "b"]


def test_coerce_unwraps_item_wrapper():
    slides = coerce_pptx_slides(
        {
            "item": [
                {"title": "A", "bullets": {"item": ["1", "2"]}},
                {"title": "B", "bullets": ["3"]},
            ]
        }
    )
    assert len(slides) == 2
    assert slides[0].bullets == ["1", "2"]
    assert slides[1].title == "B"


def test_pptx_args_uses_slides_json_string():
    args = PptxArgs.model_validate(
        {
            "template_id": "business",
            "out_path": "~/Desktop/a.pptx",
            "slides_json": '[{"title":"A","bullets":["1","2"]}]',
        }
    )
    assert args.parsed_slides()[0]["title"] == "A"
    assert args.parsed_slides()[0]["bullets"] == ["1", "2"]


def test_pptx_args_accepts_legacy_slides_and_item_wrapper():
    args = PptxArgs.model_validate(
        {
            "template_id": "business",
            "out_path": "~/Desktop/a.pptx",
            "slides": {"item": [{"title": "A", "bullets": {"item": ["1"]}}]},
        }
    )
    assert args.parsed_slides() == [{"title": "A", "bullets": ["1"]}]


def test_pptx_tool_schema_has_only_string_slide_fields():
    schema = PptxArgs.model_json_schema()
    props = schema["properties"]
    assert props["slides_json"]["anyOf"][0]["type"] == "string" or props["slides_json"].get("type") in {
        "string",
        None,
    }
    # Both slide fields must not be typed as array/object in the leaf schema.
    for key in ("slides_json", "slides"):
        node = props[key]
        # pydantic v2 optional str → anyOf [string, null]
        types = set()
        if "type" in node:
            types.add(node["type"])
        for alt in node.get("anyOf", []):
            if "type" in alt:
                types.add(alt["type"])
        assert "array" not in types
        assert types <= {"string", "null"}
